"""
Extract a .pptx file from a Drive-download JSON blob, parse slides with
python-pptx, render every slide as a PNG via PowerPoint COM automation,
and extract any embedded raster images per slide.

Outputs a small struct JSON that the model can read back to write a
verbose human-readable summary, plus per-slide PNGs and per-image files
on disk for visual analysis.

Usage:
    python extract_pptx.py <input_json_path> <work_dir> <out_struct_json_path>

Inside <work_dir> we create:
    deck.pptx
    slides/slide_001.png ... slide_NNN.png      (full slide renders)
    images/slide_NNN_img_M.<ext>                (extracted embedded images)
"""
from __future__ import annotations

import base64
import json
import sys
from pathlib import Path

from pptx import Presentation

# COM rendering is optional (only on Windows with PowerPoint installed).
try:
    import win32com.client
    import pythoncom
    HAS_PPT_COM = True
except ImportError:
    HAS_PPT_COM = False


def extract_pptx(json_blob_path: Path, out_pptx_path: Path) -> None:
    raw = json_blob_path.read_text(encoding="utf-8")
    blob = json.loads(raw)
    b64 = blob["content"]
    out_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    out_pptx_path.write_bytes(base64.b64decode(b64))


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        run_text = "".join(run.text for run in para.runs)
        if run_text.strip():
            parts.append(run_text)
    return "\n".join(parts).strip()


def extract_pictures_from_slide(slide, slide_idx: int, images_dir: Path) -> list[str]:
    """Save every raster picture shape on this slide; return relative paths."""
    paths: list[str] = []
    img_counter = 0

    def walk(shapes):
        nonlocal img_counter
        for shape in shapes:
            # Recurse into groups
            if shape.shape_type == 6:  # GROUP
                walk(shape.shapes)
                continue
            if shape.shape_type == 13:  # PICTURE
                img_counter += 1
                try:
                    img = shape.image
                    ext = img.ext or "png"
                    fname = f"slide_{slide_idx:03d}_img_{img_counter:02d}.{ext}"
                    out = images_dir / fname
                    out.write_bytes(img.blob)
                    paths.append(str(out.relative_to(images_dir.parent)))
                except Exception as e:
                    paths.append(f"[error extracting image: {e}]")

    images_dir.mkdir(parents=True, exist_ok=True)
    walk(slide.shapes)
    return paths


def slide_struct(slide, idx: int, embedded_images: list[str]) -> dict:
    title = ""
    body_blocks: list[str] = []
    image_count = 0
    table_count = 0
    chart_count = 0

    if slide.shapes.title is not None:
        t = shape_text(slide.shapes.title)
        if t:
            title = t

    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.shape_type == 13:  # PICTURE
            image_count += 1
            continue
        if shape.has_table:
            table_count += 1
            rows = []
            for row in shape.table.rows:
                rows.append(" | ".join(cell.text_frame.text.strip() for cell in row.cells))
            body_blocks.append("[TABLE]\n" + "\n".join(rows))
            continue
        if shape.has_chart:
            chart_count += 1
            body_blocks.append(f"[CHART: {shape.chart.chart_type}]")
            continue
        txt = shape_text(shape)
        if txt:
            body_blocks.append(txt)

    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    return {
        "index": idx,
        "title": title,
        "body": body_blocks,
        "image_count": image_count,
        "table_count": table_count,
        "chart_count": chart_count,
        "notes": notes,
        "embedded_images": embedded_images,
        "slide_png": None,  # filled in after COM render
    }


def render_slides_via_powerpoint(pptx_abs_path: Path, slides_dir: Path) -> list[Path]:
    """Use PowerPoint COM to export each slide as PNG. Returns list of PNG paths."""
    if not HAS_PPT_COM:
        raise RuntimeError("pywin32 not available; cannot render slides")

    slides_dir.mkdir(parents=True, exist_ok=True)
    pythoncom.CoInitialize()
    ppt_app = None
    presentation = None
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        # PowerPoint requires explicit Visible=True in some versions; use MsoTriStateTrue=-1
        # Some installs reject WithWindow=False, so we open visibly and minimize impact.
        presentation = ppt_app.Presentations.Open(
            str(pptx_abs_path), ReadOnly=True, Untitled=False, WithWindow=False
        )
        png_paths: list[Path] = []
        for i, slide in enumerate(presentation.Slides, start=1):
            out = slides_dir / f"slide_{i:03d}.png"
            # Export at 1600 wide for legibility of code/text in screenshots
            slide.Export(str(out), "PNG", 1600, 900)
            png_paths.append(out)
        return png_paths
    finally:
        if presentation is not None:
            presentation.Close()
        if ppt_app is not None:
            ppt_app.Quit()
        pythoncom.CoUninitialize()


def main() -> int:
    if len(sys.argv) != 4:
        print(__doc__)
        return 2
    json_blob_path = Path(sys.argv[1]).resolve()
    work_dir = Path(sys.argv[2]).resolve()
    out_struct_path = Path(sys.argv[3]).resolve()

    work_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = work_dir / "deck.pptx"
    slides_dir = work_dir / "slides"
    images_dir = work_dir / "images"

    extract_pptx(json_blob_path, pptx_path)

    prs = Presentation(str(pptx_path))
    slides_struct: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        embedded = extract_pictures_from_slide(slide, i, images_dir)
        slides_struct.append(slide_struct(slide, i, embedded))

    # Render slides as PNG via PowerPoint
    if HAS_PPT_COM:
        try:
            png_paths = render_slides_via_powerpoint(pptx_path, slides_dir)
            for s, p in zip(slides_struct, png_paths):
                s["slide_png"] = str(p.relative_to(work_dir.parent))
        except Exception as e:
            print(f"WARNING: slide render failed: {e}", file=sys.stderr)
    else:
        print("WARNING: pywin32 missing; skipping slide-png render", file=sys.stderr)

    out_struct_path.parent.mkdir(parents=True, exist_ok=True)
    out_struct_path.write_text(
        json.dumps({"slide_count": len(slides_struct), "slides": slides_struct}, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    print(f"OK: {len(slides_struct)} slides parsed -> {out_struct_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
